"""PowerPoint document processor for extracting text content from PPT and PPTX files.

This module provides functionality to process PowerPoint presentations and extract
text content from slides and shapes. It supports both synchronous and asynchronous
processing methods.

Example:
    >>> processor = Processor()
    >>> text = processor.process_file("presentation.pptx")
    >>> print(text)
    'Slide 1 content...\nSlide 2 content...\n'

    # Async usage with scanner
    >>> async with Scanner() as scanner:
    ...     text = await processor.process_file_async("presentation.pptx", scanner)
"""

import logging

from pptx import Presentation


class Processor:
    """PowerPoint processor for extracting text content from presentations.

    This class handles the processing of PowerPoint files (.ppt, .pptx) and
    extracts text content from all slides and shapes within the presentation.

    Attributes:
        name (str): Name identifier for the processor.
        version (str): Version number of the processor.
        supported_extensions (list): List of supported file extensions.
        logger (logging.Logger): Logger instance for the processor.

    Example:
        >>> processor = Processor()
        >>> text = processor.process_file("sample.pptx")
        >>> print(text)
        'Text from slide 1\nText from slide 2\n'
    """

    name = "PowerPoint Processor"
    version = "1.2"
    supported_extensions = [".ppt", ".pptx"]

    def __init__(self):
        """Initialize the PowerPoint processor with a configured logger."""
        self.logger = logging.getLogger("processor-powerpoint")

    def __repr__(self):
        """Return a string representation of the processor for debugging.

        Returns:
            str: A string in the format "<ClassName PowerPoint Processor v1.2>"
        """
        return f"<{self.__class__.__name__} {self.name} v{self.version}>"

    def __str__(self):
        """Return a human-readable string representation of the processor.

        Returns:
            str: A string in the format "PowerPoint Processor v1.2"
        """
        return f"{self.name} v{self.version}"

    async def process_file_async(self, file_path, scanner):
        """Process a PowerPoint file asynchronously using a scanner's executor.

        Args:
            file_path (str): Path to the PowerPoint file to process.
            scanner (Scanner): Scanner instance providing the event loop and executor.

        Returns:
            str: Extracted text content from the presentation, or empty string on error.

        Example:
            >>> async with Scanner() as scanner:
            ...     text = await processor.process_file_async("doc.pptx", scanner)
            >>> print(text)
            'Extracted text content...'
        """
        try:
            loop = scanner.loop
            text_content = await loop.run_in_executor(
                scanner.executor, self.process_file, file_path
            )
            return text_content
        except Exception as e:
            self.logger.error(f"error processing: {file_path}: {e}")
            return ""

    def process_file(self, file_path):
        """Process a PowerPoint file synchronously and extract text content.

        Extracts text from all shapes that contain text across all slides in
        the presentation.

        Args:
            file_path (str): Path to the PowerPoint file to process.

        Returns:
            str: Extracted text content from the presentation, or empty string on error.

        Example:
            >>> text = processor.process_file("presentation.pptx")
            >>> print(text)
            'Text from all slides...'
        """
        text_content = ""
        try:
            self.logger.debug(f"opening PowerPoint file: {file_path}")
            try:
                prs = Presentation(file_path)
                for slide_num, slide in enumerate(prs.slides, start=1):
                    self.logger.debug(f"processing slide: {slide_num}")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                            self.logger.debug(
                                f"extracted text from shape on slide: {slide_num}"
                            )
            except Exception as e:
                # Handle specific content-type errors that can occur with certain PowerPoint files
                if "content-type" in str(e).lower():
                    self.logger.warning(
                        f"Skipping image content-type error in {file_path}: {e}"
                    )
                    # Return any text we've extracted so far, or empty string if none
                    return text_content
                else:
                    # Re-raise if it's not the specific error we're handling
                    raise
            return text_content
        except Exception as e:
            self.logger.error(f"error processing: {file_path}: {e}")
            return ""
